import sys
import os
import zipfile
import tempfile
import collections
from pptx import Presentation
from pptx.oxml.ns import qn
from PIL import Image

def rgb_to_hex(rgb):
    try:
        return '#{0:02X}{1:02X}{2:02X}'.format(rgb[0], rgb[1], rgb[2])
    except:
        return ''

def get_dominant_colors_from_image(image_data, num_colors=2):
    """从图片字节数据中提取主色调"""
    try:
        img = Image.open(image_data)
        img = img.resize((100, 75))  # 缩小加速
        img = img.convert("RGB")
        pixels = list(img.getdata())
        counter = collections.Counter(pixels)
        most_common = counter.most_common(num_colors)
        results = []
        for color, count in most_common:
            hex_color = '#{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])
            percent = round(count / len(pixels) * 100, 1)
            results.append(f'{hex_color}({percent}%)')
        return results
    except:
        return []

def extract_images_from_pptx(pptx_path):
    """从 PPTX（ZIP）中提取所有图片，分析整体主色调"""
    all_colors = collections.Counter()
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            for name in z.namelist():
                if 'media' in name.lower() and (name.endswith('.png') or name.endswith('.jpg') or name.endswith('.jpeg')):
                    with z.open(name) as f:
                        data = f.read()
                        try:
                            img = Image.open(io.BytesIO(data))
                            img = img.resize((50, 38))
                            img = img.convert("RGB")
                            pixels = list(img.getdata())
                            all_colors.update(pixels)
                        except:
                            pass
        total = sum(all_colors.values()) or 1
        top_colors = all_colors.most_common(4)
        results = []
        for color, count in top_colors:
            hex_color = '#{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])
            percent = round(count / total * 100, 1)
            results.append(f'{hex_color}({percent}%)')
        return results
    except:
        return []

def get_background_type(slide):
    bg = slide.background
    fill = bg.fill
    try:
        if fill.type is None:
            return "无背景（白色）"
        elif fill.type == 1:
            return f"纯色背景({rgb_to_hex(fill.fore_color.rgb)})"
        elif fill.type == 2:
            return "渐变背景"
        elif fill.type == 3:
            return "图片/纹理背景"
    except:
        pass
    return "未知背景"

def collect_text_font_info(slide):
    font_colors = set()
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = ''.join(run.text for run in para.runs).strip()
                if txt:
                    texts.append(txt)
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.type is not None:
                            font_colors.add(rgb_to_hex(run.font.color.rgb))
                    except:
                        pass
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        texts.append(txt)
    return texts, font_colors

def assess_style(text_list):
    full = ' '.join(text_list)
    if any(w in full for w in ['年终', '汇报', '报告']): return '商务'
    if any(w in full for w in ['论文', '学术', '大学']): return '学术'
    if any(w in full for w in ['游戏', 'NPC', '玩法']): return '游戏/科技'
    if any(w in full for w in ['历史', '文化', '传统', '建筑']): return '中国风/历史'
    if any(w in full for w in ['课程', '教育', '培训']): return '教育'
    return '通用'

def extract_text_and_colors(filepath):
    prs = Presentation(filepath)
    all_texts = []
    bg_types = set()
    all_font_colors = set()
    result_lines = []

    # 提取 PPTX 内嵌图片的主色调
    overall_image_colors = extract_images_from_pptx(filepath)

    for i, slide in enumerate(prs.slides):
        texts, font_colors = collect_text_font_info(slide)
        all_texts.extend(texts)
        all_font_colors.update(font_colors)
        bg_type = get_background_type(slide)
        bg_types.add(bg_type)

        line = f"第{i+1}页: "
        if texts:
            line += f"文字预览: {'; '.join(texts[:3])}"
        if font_colors:
            line += f" | 字体色: {', '.join(font_colors)}"
        line += f" | 背景: {bg_type}"
        result_lines.append(line)

    overall_style = assess_style(all_texts)
    summary = f"整体风格推测: {overall_style}\n"
    summary += f"字体主色: {', '.join(all_font_colors) if all_font_colors else '未提取'}"
    if overall_image_colors:
        summary += f"\n模板整体主色调(从内嵌图片分析): {', '.join(overall_image_colors)}"
    summary += f"\n背景类型统计: {', '.join(bg_types)}\n"
    summary += '\n'.join(result_lines)
    return summary

if __name__ == '__main__':
    import io
    if len(sys.argv) < 2:
        print('')
        sys.exit(1)
    print(extract_text_and_colors(sys.argv[1]))