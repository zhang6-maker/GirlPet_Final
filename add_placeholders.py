import sys
from pptx import Presentation

def find_title_body(slide):
    """只根据字号和字数找最可能的标题和正文"""
    title_shape = None
    body_shape = None
    max_font = 0
    max_len = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        # 获取第一个run的字号
        font_size = 12
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_size = run.font.size / 12700
                break
            break
        # 标题：字号最大的
        if font_size > max_font:
            max_font = font_size
            title_shape = shape
        # 正文：字数最多的（排除标题）
        if len(text) > max_len and (title_shape is None or shape != title_shape):
            max_len = len(text)
            body_shape = shape
    return title_shape, body_shape

def replace_text_keep_format(shape, new_text):
    """清空文本并保留格式设置新文本"""
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
    else:
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        r = p.add_run()
        r.text = new_text

def add_placeholders(input_path, output_path):
    prs = Presentation(input_path)
    if not prs.slides:
        print("❌ 模板没有幻灯片")
        return
    slide = prs.slides[0]
    title_shape, body_shape = find_title_body(slide)
    if not title_shape:
        print("❌ 未找到标题")
        return
    replace_text_keep_format(title_shape, '{{title}}')
    if body_shape:
        replace_text_keep_format(body_shape, '{{content}}')
    else:
        print("⚠️ 未找到正文文本框，只在标题添加占位符")
    prs.save(output_path)
    print(f"✅ 已生成带占位符的模板: {output_path}")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("用法: python add_placeholders.py <原模板.pptx> <输出模板.pptx>")
        sys.exit(1)
    add_placeholders(sys.argv[1], sys.argv[2])