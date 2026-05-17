import sys, json
from pptx import Presentation
from pptx.util import Inches, Pt

sys.stdout.reconfigure(encoding='utf-8')

def is_title_like(shape):
    """判断形状是否像标题：字号大、字数少、位置靠上"""
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if not text or len(text) > 60:
        return False
    
    # 获取字号
    font_size = 12
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                font_size = run.font.size / 12700
                break
        break
    
    # 位置靠上（top 小于幻灯片高度的 40%）
    top = shape.top or 0
    slide_height = 6858000  # 7.5 英寸，大多数 PPT 标准高度
    
    # 字号 >= 18 且位置在上半部分，大概率是标题
    return font_size >= 18 and (top < slide_height * 0.4)

def is_body_like(shape, title_shape=None):
    """判断形状是否像正文：字数多、字号正常、不是标题"""
    if shape == title_shape:
        return False
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    return len(text) > 30

def find_title_and_body(slide):
    """返回 (标题形状, 正文形状)"""
    title_shape = None
    body_shape = None
    
    # 先找标题
    for shape in slide.shapes:
        if is_title_like(shape):
            title_shape = shape
            break
    
    # 再找正文（排除标题）
    for shape in slide.shapes:
        if is_body_like(shape, title_shape):
            body_shape = shape
            break
    
    return title_shape, body_shape

def replace_text_preserve_first_run_format(shape, new_text):
    """替换文本，保留第一个 run 的格式"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 清空所有段落
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    # 把新文本写入第一个段落
    if tf.paragraphs:
        p = tf.paragraphs[0]
        if p.runs:
            p.runs[0].text = new_text
        else:
            p.add_run().text = new_text
    else:
        p = tf.add_paragraph()
        p.add_run().text = new_text

def auto_fill_slide(slide, title_text, content_text):
    """智能填充一页幻灯片：找标题/正文文本框并替换"""
    title_shape, body_shape = find_title_and_body(slide)

    filled = False
    if title_shape:
        replace_text_preserve_first_run_format(title_shape, title_text)
        filled = True

    if body_shape:
        replace_text_preserve_first_run_format(body_shape, content_text)
        filled = True

    # 如果没找到正文框，但找到了标题框，在标题框下方添加正文段落
    if title_shape and not body_shape:
        # 在标题所在幻灯片上新增一个正文文本框
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(4)
        )
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = content_text
        p.font.size = Pt(18)
        filled = True

    return filled

def add_slide_with_text(prs, title_text, content_text):
    """追加一页新幻灯片"""
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    left, top, width, height = Inches(1), Inches(0.5), Inches(8), Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True

    left2, top2, width2, height2 = Inches(1), Inches(2), Inches(8), Inches(4)
    content_box = slide.shapes.add_textbox(left2, top2, width2, height2)
    tf2 = content_box.text_frame
    p2 = tf2.add_paragraph()
    p2.text = content_text
    p2.font.size = Pt(18)

def fill_template(template_path, output_path, json_path):
    with open(json_path, 'r', encoding='utf-8-sig') as f:
        data = json.load(f)

    title_text = data.get('title', '')
    content_text = data.get('content', '')

    prs = Presentation(template_path)

    if prs.slides:
        first_slide = prs.slides[0]
        success = auto_fill_slide(first_slide, title_text, content_text)
        if not success:
            add_slide_with_text(prs, title_text, content_text)
    else:
        add_slide_with_text(prs, title_text, content_text)

    prs.save(output_path)
    print(f"Generated: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("用法: python fill_ppt.py <模板路径> <输出路径> <JSON文件路径>")
        sys.exit(1)
    fill_template(sys.argv[1], sys.argv[2], sys.argv[3])